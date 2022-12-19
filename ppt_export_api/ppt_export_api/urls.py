"""ppt_export_api URL Configuration

The `urlpatterns` list routes URLs to views. For more information please see:
    https://docs.djangoproject.com/en/4.1/topics/http/urls/
Examples:
Function views
    1. Add an import:  from my_app import views
    2. Add a URL to urlpatterns:  path('', views.home, name='home')
Class-based views
    1. Add an import:  from other_app.views import Home
    2. Add a URL to urlpatterns:  path('', Home.as_view(), name='home')
Including another URLconf
    1. Import the include() function: from django.urls import include, path
    2. Add a URL to urlpatterns:  path('blog/', include('blog.urls'))
"""

from django.contrib import admin
from django.urls import path
from ninja import NinjaAPI
from typing import List, Any
from io import BytesIO


from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from django.http import HttpResponse


api = NinjaAPI()


@api.post("/add")
def generate_pptx(request, columns: List[str], data: List[float]) -> Any:

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # define chart data ---------------------
    chart_data = CategoryChartData()
    chart_data.categories = columns
    chart_data.add_series('Series 1', data)

    # add chart to slide --------------------
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )
    response = HttpResponse(content_type='application/vnd.ms-powerpoint')
    response['Content-Disposition'] = 'attachment; filename="sample.pptx"'
    source_stream = BytesIO()
    prs.save(source_stream)
    ppt = source_stream.getvalue()
    source_stream.close()
    response.write(ppt)
    return response



urlpatterns = [
    path("admin/", admin.site.urls),
    path("api/", api.urls),
]
